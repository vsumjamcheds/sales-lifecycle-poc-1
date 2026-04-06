"""
Build a stakeholder-style PowerPoint for the HCP Engagement POC.

Narrative order:
  1. Title
  2. The prompt (what was asked)
  3. What we understood
  4. What we developed
  5. Stack & data (compact)
  6. Placeholder slides for your screenshots (you paste images in PowerPoint)

Optional: drop PNGs into docs/ppt_assets/ and pass --embed-assets to append real screenshots
after the placeholders (same suggested titles, sorted by filename).

Run from repo root (venv activated):
  python -m src.scripts.generate_poc_deck
  python -m src.scripts.generate_poc_deck --embed-assets

Requires: python-pptx
Output: docs/HCP_Engagement_POC_Overview.pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

_ROOT = Path(__file__).resolve().parents[2]
ASSETS_DIR = _ROOT / "docs" / "ppt_assets"
OUTPUT_PATH = _ROOT / "docs" / "HCP_Engagement_POC_Overview.pptx"

BODY_PT = Pt(17)
TITLE_SUB_PT = Pt(14)

# Titles for placeholder slides (and for --embed-assets, paired with sorted PNGs)
SCREENSHOT_SUGGESTIONS = [
    "Screenshot · Territory, rep context & ranked HCP list",
    "Screenshot · Command Center — plan generation & compliance",
    "Screenshot · Strategist / agent trace (optional)",
    "Screenshot · Field note & Sync to Brain",
    "Screenshot · Conversation audit trail (per HCP)",
]


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    try:
        ph = slide.placeholders[1]
        ph.text = subtitle
        for p in ph.text_frame.paragraphs:
            p.font.size = TITLE_SUB_PT
    except (KeyError, IndexError):
        pass


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.word_wrap = True
    tf.text = bullets[0]
    tf.paragraphs[0].font.size = BODY_PT
    for line in bullets[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = BODY_PT


def _add_placeholder_visual_slide(prs: Presentation, title: str, hint: str) -> None:
    """Title + centered instruction; you replace with Paste → picture in PowerPoint."""
    layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    left, top, w, h = Inches(0.75), Inches(2.0), Inches(8.5), Inches(4.5)
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = hint
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.italic = True


def _add_image_slide(prs: Presentation, title: str, image_path: Path) -> None:
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.4), Inches(1.15), width=Inches(9.2))


def build_deck(*, embed_assets: bool) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "Intelligent HCP Engagement",
        "Proof of concept — narrative, scope, and deliverables\n(SQLite · Chroma · FastAPI · Streamlit · Claude)",
    )

    _add_bullet_slide(
        prs,
        "1 · The prompt (what we were asked)",
        [
            "Demonstrate an end-to-end field workflow for one Health Care Provider (HCP), not a slide-only story.",
            "Help the rep prioritize who to see, prepare for the call with recommendations grounded in data and prior context.",
            "Check proposed messaging against approved claims / policy before the rep acts.",
            "After the visit, capture notes, structure what matters, and write back to persistent memory.",
            "Keep an auditable trail of AI and human decisions, scoped by rep and (where relevant) by HCP.",
        ],
    )

    _add_bullet_slide(
        prs,
        "2 · What we understood",
        [
            "Territory isolation: every API call is scoped to the active rep (header), so data does not leak across regions.",
            "Two brains: structured relational data (performance, visits, commitments) plus semantic memory (vector store per HCP).",
            "Compliance is not optional copy: gatekeeper compares plan language to an approved claims corpus (verify / redline / block).",
            "Post-call capture must scrub obvious PII patterns before embedding and storage, and sync should be explicit (user-triggered).",
            "Audit is for trust: enough detail to reconstruct tool use, compliance outcomes, and human accept/reject — tied to HCP when possible.",
        ],
    )

    _add_bullet_slide(
        prs,
        "3 · What we developed",
        [
            "Scout + Capacity pipeline: rank HCPs using prescribing windows, interactions, rep capacity, and friction (seeded + API-driven).",
            "Strategist agent: Claude tool loop — fetch_hcp_performance, search_hcp_memory — returns three concrete pre-call steps (JSON).",
            "Compliance gatekeeper: embeddings + claims collection in Chroma; similarity and rules → VERIFIED / REDLINE / BLOCK with citation metadata.",
            "Scribe path: scrub → summarize (objections, tasks) → Chroma memory + SQLite interaction row + audit event.",
            "Streamlit Command Center: rep switch, live vs mock API, plan generation with spinner and plain-language pipeline copy, decision logging.",
            "Audit API: GET /api/v1/audit-logs?hcp_id=… filters JSON payload so the UI shows a per-HCP conversation trail; auto-refresh after Scribe sync.",
            "Supporting: seed script, territory isolation tests, optional Excel export, deck generator (this file).",
        ],
    )

    _add_bullet_slide(
        prs,
        "4 · Stack & persistence (where things live)",
        [
            "SQLite: reps, HCPs, prescribing signals, interactions, commitments, audit_logs.",
            "Chroma: hcp_memory (per-HCP notes) and claims_master (approved phrases for the gatekeeper).",
            "FastAPI (Uvicorn): REST API, dependency-injected DB session and rep context.",
            "Streamlit: single-page command center calling the API via httpx.",
        ],
    )

    hint = (
        "[ Insert your screenshot below this title — delete this text box after pasting ]\n\n"
        "Tip: In PowerPoint, use Insert → Pictures → This Device, "
        "or paste from the clipboard, then crop to fit."
    )
    for title in SCREENSHOT_SUGGESTIONS:
        _add_placeholder_visual_slide(prs, title, hint)

    if embed_assets:
        if not ASSETS_DIR.is_dir():
            print(f"Warning: --embed-assets but missing {ASSETS_DIR}", file=sys.stderr)
        else:
            images = sorted(ASSETS_DIR.glob("*.png"))
            for i, path in enumerate(images):
                t = SCREENSHOT_SUGGESTIONS[i] if i < len(SCREENSHOT_SUGGESTIONS) else f"Screenshot · {path.name}"
                _add_image_slide(prs, f"{t} (embedded)", path)

    _add_bullet_slide(
        prs,
        "5 · How to run the POC",
        [
            "pip install -r requirements.txt  ·  .env with ANTHROPIC_API_KEY",
            "python -m src.scripts.seed_data",
            "Terminal A: uvicorn src.api.main:app --port 8000",
            "Terminal B: streamlit run src/ui/app.py",
        ],
    )

    return prs


def main() -> int:
    ap = argparse.ArgumentParser(description="Generate HCP Engagement POC PowerPoint.")
    ap.add_argument(
        "--embed-assets",
        action="store_true",
        help=f"Append images from {ASSETS_DIR} after placeholder slides",
    )
    args = ap.parse_args()

    prs = build_deck(embed_assets=args.embed_assets)
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Wrote {OUTPUT_PATH}")
    if not args.embed_assets:
        print("Tip: add screenshots in PowerPoint on the placeholder slides, or re-run with --embed-assets")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
